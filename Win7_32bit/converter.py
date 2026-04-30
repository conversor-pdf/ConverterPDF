import os
import fitz  # PyMuPDF
from pdf2docx import Converter
import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from ebooklib import epub
import zipfile
import tempfile
import shutil

# Configuração para OCR
try:
    from pdf2image import convert_from_path
    import pytesseract
except ImportError:
    pass

def convert_to_docx(pdf_path, output_path):
    cv = Converter(pdf_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()

def convert_to_txt(pdf_path, output_path, use_ocr=False):
    if use_ocr:
        images = convert_from_path(pdf_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            for img in images:
                text = pytesseract.image_to_string(img, lang='por+eng')
                f.write(text + "\n")
    else:
        doc = fitz.open(pdf_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            for page in doc:
                f.write(page.get_text() + "\n")
        doc.close()

def convert_to_xlsx(pdf_path, output_path):
    all_data = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                all_data.extend(table)
    
    if not all_data:
        raise Exception("Nenhuma tabela encontrada no PDF.")
        
    df = pd.DataFrame(all_data[1:], columns=all_data[0] if len(all_data) > 0 else None)
    df.to_excel(output_path, index=False)

def convert_to_pptx(pdf_path, output_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    doc = fitz.open(pdf_path)
    temp_dir = tempfile.mkdtemp()
    
    try:
        for i in range(len(doc)):
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=150)
            img_path = os.path.join(temp_dir, f"page_{i}.png")
            pix.save(img_path)
            
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
        prs.save(output_path)
    finally:
        doc.close()
        shutil.rmtree(temp_dir)

def convert_to_html(pdf_path, output_path):
    doc = fitz.open(pdf_path)
    with open(output_path, 'w', encoding='utf-8') as f:
        for page in doc:
            html = page.get_text("html")
            f.write(html)
    doc.close()

def convert_to_png(pdf_path, output_path):
    doc = fitz.open(pdf_path)
    base_name = os.path.splitext(output_path)[0]
    
    if len(doc) == 1:
        page = doc.load_page(0)
        pix = page.get_pixmap(dpi=300)
        pix.save(output_path)
    else:
        # Save as ZIP
        zip_path = base_name + "_imagens.zip"
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for i in range(len(doc)):
                page = doc.load_page(i)
                pix = page.get_pixmap(dpi=300)
                img_name = f"page_{i+1}.png"
                img_path = os.path.join(tempfile.gettempdir(), img_name)
                pix.save(img_path)
                zipf.write(img_path, img_name)
                os.remove(img_path)
        output_path = zip_path # Update to reflect actual file created
        
    doc.close()
    return output_path

def convert_to_epub(pdf_path, output_path):
    book = epub.EpubBook()
    book.set_title('Documento Convertido')
    book.set_language('pt')
    
    doc = fitz.open(pdf_path)
    chapters = []
    
    for i in range(len(doc)):
        page = doc.load_page(i)
        text = page.get_text("html")
        c = epub.EpubHtml(title=f'Página {i+1}', file_name=f'page_{i+1}.xhtml')
        c.content = text
        book.add_item(c)
        chapters.append(c)
        
    doc.close()
    
    book.toc = tuple(chapters)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ['nav'] + chapters
    epub.write_epub(output_path, book)

def process_conversion(pdf_path, output_folder, formato, use_ocr=False):
    filename = os.path.basename(pdf_path)
    base_name = os.path.splitext(filename)[0]
    
    if formato == 'docx':
        out_name = f"{base_name}.docx"
        out_path = os.path.join(output_folder, out_name)
        convert_to_docx(pdf_path, out_path)
    elif formato == 'xlsx':
        out_name = f"{base_name}.xlsx"
        out_path = os.path.join(output_folder, out_name)
        convert_to_xlsx(pdf_path, out_path)
    elif formato == 'txt':
        out_name = f"{base_name}.txt"
        out_path = os.path.join(output_folder, out_name)
        convert_to_txt(pdf_path, out_path, use_ocr)
    elif formato == 'pptx':
        out_name = f"{base_name}.pptx"
        out_path = os.path.join(output_folder, out_name)
        convert_to_pptx(pdf_path, out_path)
    elif formato == 'html':
        out_name = f"{base_name}.html"
        out_path = os.path.join(output_folder, out_name)
        convert_to_html(pdf_path, out_path)
    elif formato == 'epub':
        out_name = f"{base_name}.epub"
        out_path = os.path.join(output_folder, out_name)
        convert_to_epub(pdf_path, out_path)
    elif formato == 'png':
        out_name = f"{base_name}.png"
        out_path = os.path.join(output_folder, out_name)
        out_path = convert_to_png(pdf_path, out_path)
        out_name = os.path.basename(out_path)
    else:
        raise ValueError(f"Formato não suportado: {formato}")
        
    return out_name, out_path
